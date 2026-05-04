"""Tests pour les nouveaux modules : PPTX, email, SQL retrieval, dedup, side-by-side."""

import sqlite3
import tempfile
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path

import pandas as pd
import pytest


# ── Email parsing ─────────────────────────────────────────────────────────────

class TestEmailParser:
    def test_parse_simple_eml(self, tmp_path):
        from docpipeline.parsing.email import parse_email

        msg = MIMEMultipart()
        msg["From"]    = "alice@test.com"
        msg["To"]      = "bob@test.com"
        msg["Subject"] = "Sinistre IA — dossier 2024-001"
        msg.attach(MIMEText("Bonjour Bob,\n\nVoici les détails.\n\nCordialement,\nAlice", "plain"))

        eml = tmp_path / "test.eml"
        eml.write_bytes(msg.as_bytes())

        result = parse_email(eml)
        assert result.headers["from"] == "alice@test.com"
        assert "Sinistre" in result.headers["subject"]
        assert len(result.df) >= 1
        assert "Bonjour" in result.body_text


# ── PPTX parsing ──────────────────────────────────────────────────────────────

class TestPPTXParser:
    def test_parse_real_pptx(self, tmp_path):
        from docpipeline.parsing.pptx import parse_pptx
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Synthèse 2024"
        slide.placeholders[1].text = "Premier point\nDeuxième point"

        pptx = tmp_path / "test.pptx"
        prs.save(str(pptx))

        result = parse_pptx(pptx)
        assert result.slide_count == 1
        assert "Synthèse 2024" in result.slide_titles
        assert len(result.df) >= 2


# ── SQL retrieval ─────────────────────────────────────────────────────────────

class TestSQLRetriever:
    def test_from_dataframe_creates_db(self, tmp_path):
        from docpipeline.retrieval.sql_backend import SQLRetriever

        df = pd.DataFrame({
            "page": [1, 2, 3],
            "text": [
                "La franchise est de 300 euros",
                "Le contrat couvre les dommages",
                "Garantie IA accident individuel",
            ],
        })
        db = tmp_path / "test.db"
        retriever = SQLRetriever.from_dataframe(df, db)
        assert db.exists()

        # Vérifier les tables
        with sqlite3.connect(str(db)) as conn:
            tables = [r[0] for r in conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table'"
            )]
        assert "documents" in tables
        assert "documents_fts" in tables

    def test_retrieve_finds_match(self, tmp_path):
        from docpipeline.retrieval.sql_backend import SQLRetriever

        df = pd.DataFrame({
            "page": [1, 2],
            "text": ["franchise 300 euros", "contrat assurance"],
        })
        retriever = SQLRetriever.from_dataframe(df, tmp_path / "t.db")
        result = retriever.retrieve("franchise")
        assert len(result) >= 1
        assert "franchise" in result.iloc[0]["text"].lower()


# ── Cross-doc image dedup ─────────────────────────────────────────────────────

class TestImageStore:
    def test_store_creation(self, tmp_path):
        from docpipeline.parsing.pdf.image_store import CrossDocImageStore

        store = CrossDocImageStore.open(tmp_path / "img.db", tmp_path / "imgs")
        assert (tmp_path / "img.db").exists()
        assert (tmp_path / "imgs").exists()

        stats = store.stats()
        assert stats == {"unique_images": 0, "total_occurrences": 0, "documents": 0}

    def test_ingest_real_pdf_does_not_crash(self, tmp_path):
        from docpipeline.parsing.pdf.image_store import CrossDocImageStore
        pdf = Path(__file__).parent / "fixtures" / "notice_garanties.pdf"
        if not pdf.exists():
            pytest.skip("Fixture PDF absente")

        store = CrossDocImageStore.open(tmp_path / "img.db", tmp_path / "imgs")
        result = store.ingest_pdf(pdf)
        # Le PDF de test n'a pas d'images embarquées : doit retourner zéro sans planter
        assert result == {"new": 0, "duplicates": 0}


# ── Word + PDF consolidation ──────────────────────────────────────────────────

class TestConsolidator:
    def test_consolidate_returns_unified_result(self):
        from docpipeline.parsing.word.consolidator import consolidate_word_pdf
        docx = Path(__file__).parent / "fixtures" / "contrat_assurance.docx"
        pdf  = Path(__file__).parent / "fixtures" / "notice_garanties.pdf"
        if not (docx.exists() and pdf.exists()):
            pytest.skip("Fixtures absentes")

        result = consolidate_word_pdf(docx, pdf)
        assert hasattr(result, "df")
        assert hasattr(result, "toc")
        assert hasattr(result, "annotations")
        assert isinstance(result.annotations, list)


# ── Top-level API ─────────────────────────────────────────────────────────────

class TestTopLevelAPI:
    def test_classify_via_top_api(self):
        import docpipeline
        pdf = Path(__file__).parent / "fixtures" / "notice_garanties.pdf"
        if not pdf.exists():
            pytest.skip()
        result = docpipeline.classify(pdf)
        assert result.category is not None

    def test_parse_dispatches_by_extension(self):
        import docpipeline
        pdf = Path(__file__).parent / "fixtures" / "notice_garanties.pdf"
        if not pdf.exists():
            pytest.skip()
        df = docpipeline.parse(pdf)
        assert isinstance(df, pd.DataFrame)
        assert "text" in df.columns

    def test_parse_unknown_extension_raises(self):
        import docpipeline
        with pytest.raises(ValueError, match="Format non supporté"):
            docpipeline.parse("test.xyz")

    def test_convert_unsupported_combo_raises(self, tmp_path):
        import docpipeline
        with pytest.raises(ValueError, match="Conversion non supportée"):
            docpipeline.convert(
                Path(__file__).parent / "fixtures" / "notice_garanties.pdf",
                tmp_path / "out.txt"
            )


# ── CLI ───────────────────────────────────────────────────────────────────────

class TestCLI:
    def test_version_command(self, capsys):
        from docpipeline.cli import main
        rc = main(["version"])
        assert rc == 0
        captured = capsys.readouterr()
        assert "docpipeline" in captured.out

    def test_classify_command(self, capsys):
        from docpipeline.cli import main
        pdf = Path(__file__).parent / "fixtures" / "notice_garanties.pdf"
        if not pdf.exists():
            pytest.skip()
        rc = main(["classify", str(pdf), "--json"])
        assert rc == 0
        out = capsys.readouterr().out
        assert "category" in out

    def test_help_works(self, capsys):
        from docpipeline.cli import main
        with pytest.raises(SystemExit):
            main(["--help"])
